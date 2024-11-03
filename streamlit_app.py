import streamlit as st
import pandas as pd
import fitz  # PyMuPDF for PDF processing
from sklearn.feature_extraction.text import TfidfVectorizer
import matplotlib.pyplot as plt
import seaborn as sns
from docx import Document
from bs4 import BeautifulSoup
import numpy as np
import re
from wordcloud import WordCloud
from collections import Counter
from textstat import flesch_reading_ease
from transformers import pipeline
from pptx import Presentation

# Title and file uploader on the main page
st.title("Document Analysis App")
uploaded_files = st.file_uploader("Choose files", type=["pdf", "docx", "txt", "html"], accept_multiple_files=True)

def find_palindromes(text):
    words = text.split()
    palindromes = [word for word in words if word == word[::-1] and len(word) > 2]
    return palindromes

def display_palindromes(text):
    palindromes = find_palindromes(text)
    st.subheader("Palindromic Words")
    st.write(", ".join(palindromes))
from textstat import flesch_reading_ease

def calculate_complexity(text):
    avg_sentence_length = len(text.split()) / sentence_count(text)
    avg_word_length = sum(len(word) for word in text.split()) / len(text.split())
    readability_score = flesch_reading_ease(text)
    return {
        "Average Sentence Length": avg_sentence_length,
        "Average Word Length": avg_word_length,
        "Readability Score": readability_score
    }

def display_complexity(text):
    complexity = calculate_complexity(text)
    st.subheader("Text Complexity")
    st.json(complexity)
def mood_color(score):
    if score > 5:
        return "#A9DFBF"  # light green for positive mood
    elif score < -5:
        return "#F5B7B1"  # light red for negative mood
    else:
        return "#F9E79F"  # light yellow for neutral mood

def display_text_with_mood_color(text):
    sentiment_score = simple_sentiment_analysis(text)
    color = mood_color(sentiment_score)
    st.markdown(f"<div style='background-color: {color}; padding: 10px;'>{text}</div>", unsafe_allow_html=True)

def estimate_reading_time(text, wpm=200):
    word_count = len(text.split())
    minutes = word_count / wpm
    return round(minutes, 2)

def display_reading_time(text):
    reading_time = estimate_reading_time(text)
    st.subheader("Estimated Reading Time")
    st.write(f"Approximate Reading Time: {reading_time} minutes")
def pos_tagging(text):
    doc = nlp(text)
    pos_counts = Counter([token.pos_ for token in doc])
    return pos_counts

def display_pos_tagging(text):
    st.subheader("Part-of-Speech Analysis")
    pos_counts = pos_tagging(text)
    st.bar_chart(pd.DataFrame(pos_counts.values(), index=pos_counts.keys(), columns=["Count"]))
def gendered_language_analysis(text):
    male_words = ["he", "him", "his", "man", "men"]
    female_words = ["she", "her", "hers", "woman", "women"]
    
    male_count = sum(text.lower().count(word) for word in male_words)
    female_count = sum(text.lower().count(word) for word in female_words)
    
    if male_count > female_count:
        result = "More Male-Language Oriented"
    elif female_count > male_count:
        result = "More Female-Language Oriented"
    else:
        result = "Neutral Language"
    
    return result

def display_gendered_language(text):
    st.subheader("Gendered Language Analysis")
    result = gendered_language_analysis(text)
    st.write(result)
def sentiment_by_section(text):
    paragraphs = text.split("\n\n")
    section_sentiments = [simple_sentiment_analysis(paragraph) for paragraph in paragraphs]
    plt.bar(range(len(section_sentiments)), section_sentiments, color='orange')
    plt.title("Sentiment by Section")
    plt.xlabel("Section")
    plt.ylabel("Sentiment Score")
    st.pyplot(plt)

def display_section_sentiment(text):
    st.subheader("Sentiment by Section")
    sentiment_by_section(text)
def find_jargon(text):
    common_words = set(textstat.lexicon_count(text, removepunct=True))
    technical_words = set(word for word in text.split() if len(word) > 7 and word not in common_words)
    return technical_words

def display_jargon_finder(text):
    st.subheader("Technical Jargon")
    jargon = find_jargon(text)
    st.write(", ".join(jargon))
from transformers import pipeline, T5Tokenizer

# Replace with a compatible tokenizer
tokenizer = T5Tokenizer.from_pretrained("t5-base")  # or any other supported model

# Initialize the pipeline
paraphraser = pipeline("text2text-generation", model="Vamsi/T5_Paraphrase")

# Example usage
text = "Your input text here."
paraphrased_text = paraphraser(text)
print(paraphrased_text)

def detect_tone(text):
    # Simple example, more advanced techniques would use NLP models
    if "!" in text:
        return "Casual/Excited"
    elif text.isupper():
        return "Aggressive"
    else:
        return "Neutral/Formal"
from pptx import Presentation

def create_presentation_from_text(text):
    prs = Presentation()
    slides = text.split('\n\n')  # Splitting text into slides by paragraphs
    for slide_content in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title, content = slide_content[:30], slide_content[30:]  # First 30 chars as title
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    prs.save("generated_presentation.pptx")
    return "generated_presentation.pptx"

def display_presentation_generator(text):
    st.subheader("Presentation Generator")
    presentation_path = create_presentation_from_text(text)
    st.download_button("Download Generated Presentation", presentation_path)
import fitz  # PyMuPDF for PDF manipulation

def highlight_text_in_pdf(pdf_path, keywords):
    doc = fitz.open(pdf_path)
    for page in doc:
        for keyword in keywords:
            text_instances = page.search_for(keyword)
            for inst in text_instances:
                page.add_highlight_annot(inst)
    annotated_pdf_path = "annotated_" + pdf_path
    doc.save(annotated_pdf_path)
    return annotated_pdf_path

def display_pdf_with_annotations(pdf_path, keywords):
    st.subheader("Annotated PDF")
    annotated_pdf = highlight_text_in_pdf(pdf_path, keywords)
    st.download_button("Download Annotated PDF", annotated_pdf)


def analyze_style(text):
    formal_words = set(['therefore', 'consequently', 'moreover', 'hence', 'furthermore'])
    informal_words = set(['awesome', 'cool', 'totally', 'like', 'just', 'really'])
    formal_count = sum(1 for word in text.split() if word.lower() in formal_words)
    informal_count = sum(1 for word in text.split() if word.lower() in informal_words)
    return "Formal" if formal_count > informal_count else "Informal"

def display_style_analysis(text):
    style = analyze_style(text)
    st.subheader("Writing Style Analysis")
    st.write(f"Document Style: {style}")


# Function to extract text from various formats
def extract_text(file):
    if file.type == "application/pdf":
        return extract_text_from_pdf(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_docx(file)
    elif file.type == "text/plain":
        return file.read().decode("utf-8")
    elif file.type == "text/html":
        return extract_text_from_html(file)

# PDF extraction
def extract_text_from_pdf(file):
    text = ""
    with fitz.open(stream=file.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

# DOCX extraction
def extract_text_from_docx(file):
    doc = Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

# HTML extraction
def extract_text_from_html(file):
    soup = BeautifulSoup(file.read(), 'html.parser')
    return soup.get_text()

# Text Cleaning
def clean_text(text):
    text = re.sub(r'\s+', ' ', text)  # Remove extra whitespace
    text = re.sub(r'[^\w\s]', '', text)  # Remove punctuation
    return text.strip()

# Summarization
def summarize_text(text, num_sentences=3):
    sentences = text.split('. ')
    return '. '.join(sentences[:num_sentences])

# Simple sentiment analysis
def simple_sentiment_analysis(text):
    words = text.split()
    positive_words = set(['good', 'great', 'excellent', 'positive', 'fortunate', 'correct', 'superior'])
    negative_words = set(['bad', 'poor', 'terrible', 'negative', 'unfortunate', 'wrong', 'inferior'])
    
    positive_count = sum(1 for word in words if word in positive_words)
    negative_count = sum(1 for word in words if word in negative_words)
    
    return positive_count - negative_count

# Keyword Cloud Generation
def generate_word_cloud(text):
    wordcloud = WordCloud(width=800, height=400, background_color='white').generate(text)
    plt.figure(figsize=(10, 5))
    plt.imshow(wordcloud, interpolation='bilinear')
    plt.axis('off')
    st.pyplot(plt)

# Updated Email Extraction
def extract_emails(text):
    email_pattern = r'[\w\.-]+@[\w\.-]+\.\w+'
    return re.findall(email_pattern, text)

# Link Extraction
def extract_links(text):
    return re.findall(r'https?://\S+', text)

# Frequency Distribution of Words
def word_frequency(text):
    words = text.split()
    word_counts = Counter(words)
    return word_counts.most_common(10)

# Stop Words Removal
def remove_stopwords(text):
    stop_words = set([
        'the', 'is', 'in', 'and', 'to', 'with', 'that', 'of', 'a', 'for', 'on', 
        'it', 'as', 'was', 'at', 'by', 'an', 'be', 'this', 'which', 'or', 'are'
    ])
    words = text.split()
    filtered_words = [word for word in words if word.lower() not in stop_words]
    return ' '.join(filtered_words)

# Visualizing Word Frequency
def plot_word_frequency(word_counts):
    words, counts = zip(*word_counts)
    plt.figure(figsize=(10, 5))
    sns.barplot(x=list(words), y=list(counts))
    plt.xticks(rotation=45)
    plt.title('Top 10 Most Frequent Words')
    plt.ylabel('Frequency')
    st.pyplot(plt)

# Additional Features
def character_count(text):
    return len(text)

def word_count(text):
    return len(text.split())

def sentence_count(text):
    return text.count('.') + text.count('!') + text.count('?')

def longest_word(text):
    words = text.split()
    return max(words, key=len)

def shortest_word(text):
    words = text.split()
    return min(words, key=len)

def average_word_length(text):
    words = text.split()
    return sum(len(word) for word in words) / len(words)

def top_n_words(text, n=10):
    word_counts = Counter(text.split())
    return word_counts.most_common(n)

def extract_hashtags(text):
    return re.findall(r'#\w+', text)

def extract_mentions(text):
    return re.findall(r'@\w+', text)

def top_n_sentences(text, n=3):
    sentences = text.split('. ')
    return sorted(sentences, key=len, reverse=True)[:n]

def find_unique_words(text):
    return set(text.split())

def common_phrases(text, n=5):
    phrases = re.findall(r'\b(\w+\s+\w+)\b', text)
    return Counter(phrases).most_common(n)

def keyword_extraction_tfidf(text):
    tfidf_vectorizer = TfidfVectorizer(stop_words='english')
    tfidf_matrix = tfidf_vectorizer.fit_transform([text])
    feature_names = tfidf_vectorizer.get_feature_names_out()
    dense = tfidf_matrix.todense()
    denselist = dense.tolist()
    return pd.DataFrame(denselist, columns=feature_names).T.sort_values(0, ascending=False).head(10)

def count_paragraphs(text):
    return text.count('\n') + 1

def check_spelling(text):
    from spellchecker import SpellChecker
    spell = SpellChecker()
    words = text.split()
    misspelled = spell.unknown(words)
    return list(misspelled)

def extract_dates(text):
    return re.findall(r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b', text)

def extract_time(text):
    return re.findall(r'\b\d{1,2}:\d{2}\s*[AP]?[M]?\b', text)

def text_to_uppercase(text):
    return text.upper()

def text_to_lowercase(text):
    return text.lower()

def text_title_case(text):
    return text.title()

def remove_special_characters(text):
    return re.sub(r'[^\w\s]', '', text)

def extract_citations(text):
    return re.findall(r'\([^\)]+\)', text)

def extract_references(text):
    return re.findall(r'\[.*?\]', text)

def count_links(text):
    return len(extract_links(text))

def count_emails(text):
    return len(extract_emails(text))

def extract_quotes(text):
    return re.findall(r'"(.*?)"', text)

def find_most_common_char(text):
    return Counter(text).most_common(1)

def text_statistics(text):
    return {
        "Character Count": character_count(text),
        "Word Count": word_count(text),
        "Sentence Count": sentence_count(text),
        "Longest Word": longest_word(text),
        "Shortest Word": shortest_word(text),
        "Average Word Length": average_word_length(text)
    }

# Main application logic
if uploaded_files:
    all_texts = []
    for uploaded_file in uploaded_files:
        text = extract_text(uploaded_file)
        all_texts.append(text)

    full_text = "\n\n".join(all_texts)
    
    # Text Analysis Section
    st.subheader("Extracted Text")
    st.write(full_text)

    # Cleaned Text
    cleaned_text = clean_text(full_text)
    st.subheader("Cleaned Text")
    st.write(cleaned_text)

    # Text Statistics
    stats = text_statistics(cleaned_text)
    st.subheader("Text Statistics")
    st.json(stats)

    # Simple Sentiment Analysis
    sentiment_score = simple_sentiment_analysis(cleaned_text)
    st.subheader("Sentiment Analysis")
    st.write(f"Sentiment Score: {sentiment_score}")

    # Summary of Text
    st.subheader("Text Summary")
    summary = summarize_text(cleaned_text)
    st.write(summary)

    # Keyword Extraction using TF-IDF
    st.subheader("Keyword Extraction")
    df_tfidf = keyword_extraction_tfidf(cleaned_text)
    st.write(df_tfidf)

    # Generate Word Cloud
    st.subheader("Keyword Cloud")
    generate_word_cloud(cleaned_text)

    # Email and Link Extraction
    st.subheader("Extracted Emails")
    emails = extract_emails(cleaned_text)
    st.write(emails)
    st.write(f"Total Emails: {count_emails(cleaned_text)}")

    st.subheader("Extracted Links")
    links = extract_links(cleaned_text)
    st.write(links)
    st.write(f"Total Links: {count_links(cleaned_text)}")

    st.subheader("Extracted Dates")
    dates = extract_dates(cleaned_text)
    st.write(dates)

    # Frequency Distribution
    st.subheader("Top 10 Most Frequent Words")
    freq_words = word_frequency(cleaned_text)
    plot_word_frequency(freq_words)

    # Unique Words
    st.subheader("Unique Words")
    st.write(find_unique_words(cleaned_text))

    # Quotes
    st.subheader("Extracted Quotes")
    quotes = extract_quotes(cleaned_text)
    st.write(quotes)

    # Check for spelling errors
    st.subheader("Spelling Errors")
    spelling_errors = check_spelling(cleaned_text)
    st.write(spelling_errors)

    # Paragraph Count
    st.subheader("Paragraph Count")
    st.write(count_paragraphs(cleaned_text))

    # Citations
    st.subheader("Extracted Citations")
    citations = extract_citations(cleaned_text)
    st.write(citations)

    # References
    st.subheader("Extracted References")
    references = extract_references(cleaned_text)
    st.write(references)

    # Special Character Removal
    st.subheader("Text without Special Characters")
    st.write(remove_special_characters(cleaned_text))

    # Download processed text if needed
    st.sidebar.subheader("Download Processed Text")
    if st.sidebar.button("Download"):
        st.sidebar.download_button("Download Cleaned Text", cleaned_text)

# Sidebar Title
st.sidebar.text("Document Analysis App")
